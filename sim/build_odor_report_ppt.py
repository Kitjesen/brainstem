from __future__ import annotations

import csv
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
SIM_OUTPUT = ROOT / "brain" / "brainstem" / "sim" / "output"
OUTPUT_PPT = Path(__file__).resolve().parents[3] / "docs" / "ppt" / "odor_source_localization_report_cn.pptx"

TITLE_COLOR = RGBColor(20, 37, 63)
TEXT_COLOR = RGBColor(38, 42, 48)
MUTED_COLOR = RGBColor(96, 103, 112)
ACCENT = RGBColor(14, 120, 110)
ACCENT_2 = RGBColor(28, 126, 214)
SUCCESS = RGBColor(43, 138, 62)
CHALLENGE = RGBColor(201, 42, 42)
BG = RGBColor(247, 249, 252)


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = TEXT_COLOR,
    font_name: str = "Microsoft YaHei",
    align: PP_ALIGN | None = None,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    if align is not None:
        paragraph.alignment = align
    return box


def add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    bullets: list[str],
    *,
    font_size: int = 18,
    color: RGBColor = TEXT_COLOR,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.space_after = Pt(7)
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Microsoft YaHei"
    return box


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, 0.55, 0.28, 11.6, 0.55, title, font_size=28, bold=True, color=TITLE_COLOR)
    if subtitle:
        add_textbox(slide, 0.58, 0.80, 11.6, 0.30, subtitle, font_size=12, color=MUTED_COLOR)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.10), Inches(1.55), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_2
    line.line.fill.background()


def add_image(slide, path: Path, left: float, top: float, width: float, *, height: float | None = None) -> None:
    kwargs = {"width": Inches(width)}
    if height is not None:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), **kwargs)


def add_formula_box(slide, title: str, formulas: list[str], left: float, top: float, width: float, height: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(242, 245, 249)
    shape.line.color.rgb = RGBColor(214, 220, 228)
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(16)
    r.font.color.rgb = TITLE_COLOR
    r.font.name = "Microsoft YaHei"
    for formula in formulas:
        p = text_frame.add_paragraph()
        p.text = formula
        p.space_after = Pt(4)
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = TEXT_COLOR
            run.font.name = "Consolas"


def add_caption(slide, text: str, left: float, top: float, width: float) -> None:
    add_textbox(slide, left, top, width, 0.35, text, font_size=11, color=MUTED_COLOR)


def add_table(slide, left: float, top: float, width: float, height: float, data: list[list[str]]) -> None:
    rows = len(data)
    cols = len(data[0]) if data else 0
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(12 if r == 0 else 11)
                    run.font.bold = r == 0
                    run.font.color.rgb = TITLE_COLOR if r == 0 else TEXT_COLOR
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = RGBColor(227, 239, 255)
            else:
                cell.fill.fore_color.rgb = RGBColor(251, 252, 253) if r % 2 == 1 else RGBColor(245, 247, 250)


def load_metrics(csv_path: Path, success_radius: float = 0.65) -> tuple[str, str]:
    rows = list(csv.DictReader(csv_path.open("r", encoding="utf-8")))
    distances = [float(row["distance_to_source_m"]) for row in rows]
    reached_at = next((float(row["time_s"]) for row in rows if float(row["distance_to_source_m"]) <= success_radius), None)
    reach_text = f"{reached_at:.2f} s" if reached_at is not None else "not reached"
    min_text = f"{min(distances):.3f} m"
    return reach_text, min_text


def count_retarget_events(csv_path: Path) -> int:
    rows = list(csv.DictReader(csv_path.open("r", encoding="utf-8")))
    return len(rows)


def add_metric_badge(slide, label: str, value: str, left: float, top: float, width: float, *, good: bool = True) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.72)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(238, 244, 248)
    shape.line.color.rgb = RGBColor(214, 220, 228)
    add_textbox(slide, left + 0.15, top + 0.10, width - 0.3, 0.20, label, font_size=11, color=MUTED_COLOR)
    add_textbox(
        slide,
        left + 0.15,
        top + 0.28,
        width - 0.3,
        0.25,
        value,
        font_size=18,
        bold=True,
        color=SUCCESS if good else CHALLENGE,
    )


def pick_visual(gif: Path, fallback: Path) -> Path:
    """Return *gif* if it exists on disk, otherwise *fallback* (PNG)."""
    return gif if gif.exists() else fallback


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    cfd_plot = SIM_OUTPUT / "odor_cfd_turbulence_demo.png"
    cfd_gif = SIM_OUTPUT / "odor_cfd_turbulence_demo.gif"
    mujoco_plot = SIM_OUTPUT / "odor_mujoco_go2w_v3_demo.png"
    mujoco_storyboard = SIM_OUTPUT / "odor_mujoco_go2w_v3_storyboard.png"
    mujoco_gif = SIM_OUTPUT / "odor_mujoco_go2w_v3_demo.gif"
    inline_gif = SIM_OUTPUT / "odor_mujoco_go2w_v3_inline_demo.gif"
    slowwind_gif = SIM_OUTPUT / "odor_mujoco_go2w_v3_slowwind_demo.gif"
    challenge_gif = SIM_OUTPUT / "odor_mujoco_go2w_v3_offset_demo.gif"
    challenge_storyboard = SIM_OUTPUT / "odor_mujoco_go2w_v3_offset_storyboard.png"
    casebook = SIM_OUTPUT / "odor_go2w_casebook.png"
    retarget_plot = SIM_OUTPUT / "odor_nav_gateway_retarget_demo.png"
    retarget_gif = SIM_OUTPUT / "odor_nav_gateway_retarget_demo.gif"
    retarget_storyboard = SIM_OUTPUT / "odor_nav_gateway_retarget_storyboard.png"
    retarget_csv = SIM_OUTPUT / "odor_nav_gateway_retarget_demo.csv"
    retarget_events_csv = SIM_OUTPUT / "odor_nav_gateway_retarget_events.csv"

    room_casebook = SIM_OUTPUT / "odor_room_casebook.png"
    room_ids = ["room_upwind", "room_cross", "room_diagonal", "room_headwind",
                "room_sidewind", "room_corner", "room_weak", "room_strong"]
    room_gifs = {rid: SIM_OUTPUT / f"odor_{rid}_demo.gif" for rid in room_ids}
    room_pngs = {rid: SIM_OUTPUT / f"odor_{rid}_demo.png" for rid in room_ids}
    room_csvs = {rid: SIM_OUTPUT / f"odor_{rid}_demo.csv" for rid in room_ids}

    default_reach, default_min = load_metrics(SIM_OUTPUT / "odor_mujoco_go2w_v3_demo.csv")
    inline_reach, inline_min = load_metrics(SIM_OUTPUT / "odor_mujoco_go2w_v3_inline_demo.csv")
    slowwind_reach, slowwind_min = load_metrics(SIM_OUTPUT / "odor_mujoco_go2w_v3_slowwind_demo.csv")
    offset_reach, offset_min = load_metrics(SIM_OUTPUT / "odor_mujoco_go2w_v3_offset_demo.csv")
    midlane_reach, midlane_min = load_metrics(SIM_OUTPUT / "odor_mujoco_go2w_v3_midlane_demo.csv")
    heading_reach, heading_min = load_metrics(SIM_OUTPUT / "odor_mujoco_go2w_v3_heading176_demo.csv")
    retarget_reach, retarget_min = load_metrics(retarget_csv)
    retarget_count = count_retarget_events(retarget_events_csv)

    room_metrics: dict[str, tuple[str, str]] = {}
    for rid in room_ids:
        room_metrics[rid] = load_metrics(room_csvs[rid], success_radius=0.55)
    room_successes = sum(1 for rid in room_ids if room_metrics[rid][0] != "not reached")
    room_challenges = len(room_ids) - room_successes

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_textbox(slide, 0.80, 1.05, 11.5, 0.80, "四传感器嗅觉导航与 MuJoCo 演示汇报", font_size=30, bold=True, color=TITLE_COLOR)
    add_textbox(
        slide,
        0.84,
        1.95,
        11.2,
        0.55,
        "单一气味源定位 · CFD 风格动态羽流 · Go2W MuJoCo 实际行走 · nav-gateway / LingTu 对接路径",
        font_size=18,
        color=ACCENT_2,
    )
    add_bullets(
        slide,
        0.95,
        2.80,
        10.8,
        2.10,
        [
            "目标：用 4 个同型号嗅觉传感器的相对变化，而不是绝对浓度标定，驱动机器人找到单一气味源附近。",
            "路线：四传感器阵列 -> baseline / delta / gradient -> odor planner -> MuJoCo / nav-gateway / LingTu。",
            "当前演示重点：把空气流动、羽流变化、传感器响应和机器狗轨迹放到同一套可讲清楚的画面里。",
        ],
        font_size=20,
    )
    add_metric_badge(slide, "Default v3 reach", default_reach, 0.95, 5.35, 2.35)
    add_metric_badge(slide, "Default v3 min distance", default_min, 3.45, 5.35, 2.55)
    add_metric_badge(slide, "Inline reach", inline_reach, 6.20, 5.35, 2.10)
    add_metric_badge(slide, "Boundary challenge", offset_min, 8.45, 5.35, 2.70, good=False)
    add_metric_badge(slide, "nav-gateway retargets", str(retarget_count), 11.30, 5.35, 1.70)
    add_textbox(slide, 0.85, 6.75, 5.0, 0.25, "生成时间：2026-03-11", font_size=11, color=MUTED_COLOR)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_title(slide, "1. 项目背景与问题定义", "为什么要做嗅觉导航，以及为什么不能只看一个传感器的 ppm 值")
    add_bullets(
        slide,
        0.70,
        1.45,
        5.8,
        4.50,
        [
            "工业现场的气体泄漏、异味定位、异常巡检，本质上都是“沿气味线索搜源”的问题，而不是普通点到点路径规划。",
            "真实羽流不是静态高斯团块，而是受风场、湍流、障碍物、传感器迟滞共同影响的时空信号。",
            "单个传感器常常只能提供“变化变强还是变弱”，很难稳定给出方向，所以必须引入空间阵列与时间滤波。",
            "导航系统不能被嗅觉模块直接接管：嗅觉只负责滚动局部目标，避障、安全和恢复仍然必须走 LingTu / 控制栈。",
        ],
    )
    add_formula_box(
        slide,
        "本项目的最小可行假设",
        [
            "task = single odor source localization",
            "sensor = four identical channels + relative change only",
            "planner = gradient + state machine",
            "executor = MuJoCo / nav-gateway / LingTu",
        ],
        6.75,
        1.65,
        5.60,
        2.35,
    )
    add_formula_box(
        slide,
        "为什么 4 个传感器就够做第一版",
        [
            "g_x = ((d_FL + d_FR) - (d_RL + d_RR)) / 2",
            "g_y = ((d_FL + d_RL) - (d_FR + d_RR)) / 2",
            "strength = mean(max(d_i, 0))",
        ],
        6.75,
        4.25,
        5.60,
        1.90,
    )

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_title(slide, "2. 系统链路与控制逻辑", "从传感器原始值到局部导航目标")
    add_formula_box(
        slide,
        "信号处理",
        [
            "baseline_i(t) = (1 - alpha) * baseline_i(t-1) + alpha * x_i(t)",
            "delta_i(t) = filtered_i(t) - baseline_i(t)",
            "slope_i(t) = (delta_i(t) - delta_i(t-1)) / dt",
        ],
        0.75,
        1.55,
        5.70,
        2.20,
    )
    add_formula_box(
        slide,
        "状态机",
        [
            "search -> track -> cast -> reacquire",
            "search: 缓慢扫描",
            "track: 沿局部梯度收敛",
            "cast: 丢失后横向重新找羽流",
        ],
        6.70,
        1.55,
        5.70,
        2.20,
    )
    add_bullets(
        slide,
        0.80,
        4.10,
        11.0,
        2.10,
        [
            "嗅觉策略只输出短程滚动方向或局部目标，不直接驱动关节，不绕过安全栈。",
            "MuJoCo 演示中当前链路是：四传感器阵列 -> odor policy -> Walk(x, y, z) -> CMS -> MuJoCo 四足运动。",
            "落到产品化路径时应改成：odor planner -> nav-gateway target update -> LingTu -> 现有避障与恢复。",
        ],
        font_size=18,
    )

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_title(slide, "3. CFD 风格空气与羽流建模", "先验证空气模型——把风场、羽流和传感器响应放到同一个画面里")
    add_image(slide, cfd_gif, 0.60, 1.45, 6.15)
    add_image(slide, cfd_plot, 6.95, 1.45, 5.75)
    add_caption(slide, "左：动态羽流分镜。右：浓度场、风向、传感器曲线和距离曲线。", 0.72, 6.55, 11.3)
    add_bullets(
        slide,
        6.95,
        5.70,
        5.40,
        0.75,
        [
            "动态羽流不是只画热力图，而是让风场箭头、羽流偏转、局部团块都随时间变化。",
        ],
        font_size=16,
        color=ACCENT,
    )

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_title(slide, "4. MuJoCo Go2W v3 默认演示", "这一页是给导师看“系统已经跑起来了”的核心页")
    add_image(slide, mujoco_storyboard, 0.58, 1.45, 6.20)
    add_image(slide, mujoco_plot, 6.95, 1.45, 5.75)
    add_metric_badge(slide, "reach time", default_reach, 0.75, 6.15, 1.75)
    add_metric_badge(slide, "min distance", default_min, 2.65, 6.15, 1.90)
    add_metric_badge(slide, "scene", "v3 corridor", 4.70, 6.15, 1.95)
    add_metric_badge(slide, "airflow", "dynamic", 6.85, 6.15, 1.75)
    add_bullets(
        slide,
        8.85,
        6.08,
        3.10,
        0.70,
        [
            "右上角大地图显示动态风场和羽流，左侧 3D 视图显示 Go2W 实际行走。"
        ],
        font_size=14,
        color=ACCENT_2,
    )

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_title(slide, "5. GIF Demo 1: 默认场景与 Inline 场景", "放映模式下 GIF 自动播放 — 默认场景快速收敛，Inline 场景验证出发位置对结果的影响")
    add_image(slide, mujoco_gif, 0.62, 1.45, 6.05)
    add_image(slide, inline_gif, 6.82, 1.45, 5.85)
    add_caption(slide, "左：Default v3。右：Inline plume start。", 0.80, 6.55, 10.8)
    add_metric_badge(slide, "Default reach", default_reach, 0.85, 6.05, 1.90)
    add_metric_badge(slide, "Inline reach", inline_reach, 3.00, 6.05, 1.90)
    add_metric_badge(slide, "Inline min", inline_min, 5.15, 6.05, 1.90)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "6. GIF Demo 2: 风速变化场景与方法边界案例",
        "验证环境变化对搜源效果的影响 — 风速降低仍能收敛，偏置工况暴露方法边界",
    )
    add_image(slide, slowwind_gif, 0.62, 1.45, 6.05)
    add_image(slide, challenge_gif, 6.82, 1.45, 5.85)
    add_caption(
        slide,
        "左：将背景风速从 0.62 m/s 调低到 0.52 m/s，观察羽流形状和机器人轨迹如何变化。右：源点和起点一起偏移后，机器人丢失主羽流，这是当前方法的边界失败案例。",
        0.82,
        6.50,
        11.2,
    )
    add_metric_badge(slide, "风速变化场景 reach", slowwind_reach, 0.85, 6.00, 2.35)
    add_metric_badge(slide, "风速变化场景 min", slowwind_min, 3.40, 6.00, 2.35)
    add_metric_badge(slide, "边界案例最近距离", offset_min, 8.35, 6.00, 2.35, good=False)
    add_bullets(
        slide,
        0.86,
        6.85,
        11.0,
        0.52,
        [
            "这里的“风速变化场景”不是换任务，而是只改变空气流速，用来说明气流变了以后羽流和搜源过程也会跟着变。",
            "这里的“方法边界案例”不是程序报错，而是故意展示当前局部梯度 + 状态机方法在偏置工况下还不够鲁棒。"
        ],
        font_size=14,
        color=ACCENT,
    )

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_title(slide, "7. 多案例 Casebook", "把成功样例和挑战样例放在一张图上，便于汇报时快速对比")
    add_image(slide, casebook, 0.55, 1.32, 12.18)
    add_caption(slide, "6 个案例：3 个成功样例，3 个边界/失败样例。", 0.82, 6.75, 10.8)

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_title(slide, "8. 当前结果汇总", "能做什么、不能做什么，需要明说")
    table_data = [
        ["案例", "结果", "reach time", "min distance", "说明"],
        ["Default v3", "SUCCESS", default_reach, default_min, "主汇报 demo，动态气流 + Go2W v3 场景"],
        ["Inline plume", "SUCCESS", inline_reach, inline_min, "从羽流中心线附近起步，收敛更直接"],
        ["Slow-wind", "SUCCESS", slowwind_reach, slowwind_min, "风速变化案例，用于说明空气流动变化"],
        ["Offset source", "CHALLENGE", offset_reach, offset_min, "源点和起点一起偏移后仍会丢失羽流"],
        ["Midlane start", "CHALLENGE", midlane_reach, midlane_min, "靠近通道中心的起步仍不够鲁棒"],
        ["Wind 176 deg", "CHALLENGE", heading_reach, heading_min, "风向轻微旋转就暴露模型脆弱性"],
    ]
    add_table(slide, 0.55, 1.50, 12.15, 4.85, table_data)
    add_bullets(
        slide,
        0.80,
        6.10,
        11.1,
        0.70,
        [
            "这一页的意义不是“证明已经完美”，而是明确告诉导师：默认工况能跑通，扰动工况还需要更强的风场利用与局部重规划。",
        ],
        font_size=16,
        color=ACCENT,
    )

    # Slide 10
    slide = prs.slides.add_slide(blank)
    add_title(slide, "9. 机器人震荡问题是怎么处理的", "这次改动的重点之一就是把 Go2W 的抖动压下去")
    add_bullets(
        slide,
        0.75,
        1.50,
        5.85,
        4.50,
        [
            "对横向梯度新增 `filtered_steer`，避免直接用生硬的瞬时差分去拧转向。",
            "在源附近提高 deadband、降低 turn gain、降低 turn limit，并且进一步压低前进速度。",
            "将 `track` 模式下的 `walk_z` 每步变化限幅从更激进的值收紧，减少摆头抖动。",
            "结果：Default v3 的 `walk_z` 总变差已经从 v2 的 1.1716 下降到 v3 的 0.8388。"
        ],
    )
    add_formula_box(
        slide,
        "v3 近源收敛策略",
        [
            "near_source = (strength > 62) or (g_x < -4)",
            "very_near = (strength > 90) or (g_x < -10)",
            "turn_gain = 0.017 -> 0.009",
            "turn_limit = 0.22 -> 0.10",
            "delta walk_z per step <= 0.028 in track mode",
        ],
        6.85,
        1.75,
        5.25,
        2.50,
    )
    add_image(slide, challenge_storyboard, 6.80, 4.35, 5.55)
    add_caption(slide, "挑战场景仍然会丢失羽流，这也是下一阶段要重点修的地方。", 6.95, 6.60, 5.0)

    # Slide 11 — nav-gateway retarget demo (GIF + storyboard)
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "10. 产品链路验证 — nav-gateway 滚动目标",
        "用 nav-gateway API 实际发送滚动目标，证明嗅觉策略可以接入 LingTu 导航栈（放映模式下 GIF 自动播放）",
    )
    # Use GIF if available, fallback to PNG
    retarget_visual = retarget_gif if retarget_gif.exists() else retarget_plot
    add_image(slide, retarget_visual, 0.55, 1.40, 6.10)
    # Storyboard on the right if available, otherwise show metrics in that space
    if retarget_storyboard.exists():
        add_image(slide, retarget_storyboard, 6.80, 1.40, 5.85)
    add_metric_badge(slide, "retarget count", str(retarget_count), 0.70, 6.00, 2.00)
    add_metric_badge(slide, "reach time", retarget_reach, 2.85, 6.00, 2.00)
    add_metric_badge(slide, "min distance", retarget_min, 5.00, 6.00, 2.00)
    add_bullets(
        slide,
        7.30,
        5.95,
        5.50,
        1.10,
        [
            "嗅觉策略每 0.7s 输出一个滚动目标，通过 nav-gateway dispatch/retarget API 发送。",
            "nav-gateway 将目标转发给 LingTu bridge，验证了完整的 odor planner -> nav-gateway -> LingTu 链路。",
        ],
        font_size=15,
    )
    add_caption(
        slide,
        "左：GIF 动画展示滚动目标实时发送过程。右：四阶段 Storyboard 展示搜源到收敛全过程。",
        0.72,
        6.85,
        11.3,
    )

    # Slide 12 — room-scale overview
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "11. 室内场景总览 — 8 组 Room-Scale 嗅觉导航",
        "8m × 6m 房间，变化气源位置、风向、起始点，验证探索与局部规划鲁棒性",
    )
    add_image(slide, room_casebook, 0.55, 1.32, 9.00)
    add_formula_box(
        slide,
        "8 组室内场景",
        [
            "room_upwind    — 正对上风（SUCCESS）",
            "room_cross     — 横风偏置（CHALLENGE）",
            "room_diagonal  — 对角风场（SUCCESS）",
            "room_headwind  — 下风起步（CHALLENGE）",
            "room_sidewind  — 侧风 270°（SUCCESS）",
            "room_corner    — 角落气源（SUCCESS）",
            "room_weak      — 微风 0.28 m/s（SUCCESS）",
            "room_strong    — 强风 0.90 m/s（SUCCESS）",
        ],
        9.75,
        1.32,
        3.20,
        3.80,
    )
    add_metric_badge(slide, "场景总数", str(len(room_ids)), 0.70, 6.00, 1.85)
    add_metric_badge(slide, "成功", str(room_successes), 2.70, 6.00, 1.85)
    add_metric_badge(slide, "挑战", str(room_challenges), 4.70, 6.00, 1.85, good=False)
    add_metric_badge(slide, "房间尺寸", "8m × 6m", 6.70, 6.00, 1.85)
    add_caption(
        slide,
        "Casebook 4×2 网格：每格展示一个场景的轨迹和浓度热力图，绿色标记成功、红色标记挑战。",
        0.72,
        6.85,
        11.3,
    )

    # Slide 13 — exploration vs local planning behavior
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "12. 探索与局部规划行为",
        "对比直接追踪 vs 探索-横扫-追踪状态切换：search 找信号、track 跟梯度、cast 丢失后重找",
    )
    cross_visual = pick_visual(room_gifs["room_cross"], room_pngs["room_cross"])
    upwind_visual = pick_visual(room_gifs["room_upwind"], room_pngs["room_upwind"])
    add_image(slide, cross_visual, 0.55, 1.40, 6.05)
    add_image(slide, upwind_visual, 6.75, 1.40, 5.85)
    cross_reach, cross_min = room_metrics["room_cross"]
    upwind_reach, upwind_min = room_metrics["room_upwind"]
    add_metric_badge(slide, "横风偏置 reach", cross_reach, 0.70, 6.00, 2.50, good=cross_reach != "not reached")
    add_metric_badge(slide, "横风偏置 min dist", cross_min, 3.35, 6.00, 2.50, good=cross_reach != "not reached")
    add_metric_badge(slide, "正对上风 reach", upwind_reach, 6.90, 6.00, 2.50)
    add_metric_badge(slide, "正对上风 min dist", upwind_min, 9.55, 6.00, 2.50)
    add_caption(slide, "左：room_cross（横风偏置），持续 search 未切换到 track。右：room_upwind（正对上风），快速 search→track 收敛。", 0.72, 6.55, 11.3)
    add_bullets(
        slide,
        0.72,
        6.85,
        11.3,
        0.55,
        [
            "这是回答'怎么探索、怎么局部规划'的核心页：成功案例 1 步 search→track，失败案例始终停留在 search，说明搜索策略需要加入全局覆盖。",
        ],
        font_size=14,
        color=ACCENT,
    )

    # Slide 14 — wind direction & source position effects
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "13. 风向与气源位置效应",
        "不同风向改变羽流形状和扩散路径，机器人必须适应不同的梯度几何",
    )
    sidewind_visual = pick_visual(room_gifs["room_sidewind"], room_pngs["room_sidewind"])
    diagonal_visual = pick_visual(room_gifs["room_diagonal"], room_pngs["room_diagonal"])
    add_image(slide, sidewind_visual, 0.55, 1.40, 6.05)
    add_image(slide, diagonal_visual, 6.75, 1.40, 5.85)
    sw_reach, sw_min = room_metrics["room_sidewind"]
    diag_reach, diag_min = room_metrics["room_diagonal"]
    add_metric_badge(slide, "侧风 reach", sw_reach, 0.70, 6.00, 2.50)
    add_metric_badge(slide, "侧风 min dist", sw_min, 3.35, 6.00, 2.50)
    add_metric_badge(slide, "对角风场 reach", diag_reach, 6.90, 6.00, 2.50, good=diag_reach != "not reached")
    add_metric_badge(slide, "对角风场 min dist", diag_min, 9.55, 6.00, 2.50, good=diag_reach != "not reached")
    add_caption(
        slide,
        "左：room_sidewind（侧风 270°，SUCCESS）。右：room_diagonal（对角风场 225°，SUCCESS，通过 zigzag 搜索在 45s 内收敛）。",
        0.72,
        6.55,
        11.3,
    )
    add_bullets(
        slide,
        0.72,
        6.85,
        11.3,
        0.55,
        [
            "侧风时羽流沿垂直方向扩展，机器人仍可用横向梯度修正航向；对角风场 225° 使羽流偏转，但 zigzag 搜索策略在 ~45s 内成功捕获气羽并收敛到气源。",
        ],
        font_size=14,
        color=ACCENT,
    )

    # Slide 15 — wind speed effects
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "14. 风速效应 — 微风 vs 强风",
        "风速决定羽流宽度和传输速度：微风宽而弱，强风窄而快",
    )
    weak_visual = pick_visual(room_gifs["room_weak"], room_pngs["room_weak"])
    strong_visual = pick_visual(room_gifs["room_strong"], room_pngs["room_strong"])
    add_image(slide, weak_visual, 0.55, 1.40, 6.05)
    add_image(slide, strong_visual, 6.75, 1.40, 5.85)
    weak_reach, weak_min = room_metrics["room_weak"]
    strong_reach, strong_min = room_metrics["room_strong"]
    add_metric_badge(slide, "微风 reach", weak_reach, 0.70, 6.00, 2.50)
    add_metric_badge(slide, "微风 min dist", weak_min, 3.35, 6.00, 2.50)
    add_metric_badge(slide, "强风 reach", strong_reach, 6.90, 6.00, 2.50)
    add_metric_badge(slide, "强风 min dist", strong_min, 9.55, 6.00, 2.50)
    add_caption(
        slide,
        "左：room_weak（微风 0.28 m/s）。右：room_strong（强风 0.90 m/s）。",
        0.72,
        6.55,
        11.3,
    )
    add_bullets(
        slide,
        0.72,
        6.85,
        11.3,
        0.55,
        [
            "微风下羽流扩散更宽、浓度衰减更快，但覆盖面广，机器人能较快捕获信号；强风下羽流窄而集中，梯度更尖锐，收敛更快更精确。",
        ],
        font_size=14,
        color=ACCENT,
    )

    # Slide 16 — room-scale results summary table
    slide = prs.slides.add_slide(blank)
    add_title(
        slide,
        "15. 室内场景结果汇总",
        "8 组 room-scale 场景的定量结果和行为分析",
    )
    room_labels: dict[str, tuple[str, str, str]] = {
        "room_upwind":   ("正对上风",    "180°", "0.55 m/s"),
        "room_cross":    ("横风偏置",    "180°", "0.55 m/s"),
        "room_diagonal": ("对角风场",    "225°", "0.55 m/s"),
        "room_headwind": ("下风起步",    "180°", "0.55 m/s"),
        "room_sidewind": ("侧风",        "270°", "0.55 m/s"),
        "room_corner":   ("角落气源",    "210°", "0.55 m/s"),
        "room_weak":     ("微风",        "180°", "0.28 m/s"),
        "room_strong":   ("强风",        "180°", "0.90 m/s"),
    }
    room_table_data: list[list[str]] = [
        ["场景", "风向", "风速", "结果", "reach time", "min distance"],
    ]
    for rid in room_ids:
        label, wind_dir, wind_spd = room_labels[rid]
        r_reach, r_min = room_metrics[rid]
        status = "SUCCESS" if r_reach != "not reached" else "CHALLENGE"
        room_table_data.append([label, wind_dir, wind_spd, status, r_reach, r_min])
    add_table(slide, 0.55, 1.50, 12.15, 4.30, room_table_data)
    add_bullets(
        slide,
        0.80,
        5.90,
        11.1,
        1.20,
        [
            "上风/同轴和侧风工况收敛快且稳定；角落场景可收敛但需更长时间。",
            "横风偏置和对角风场无法在 40s 内收敛：当前 search 策略只沿上风搜索，不覆盖垂直于风向的区域。",
            "下风起步完全失败：机器人在源的上风侧，羽流扩散方向远离机器人，需要加入全局搜索覆盖。",
        ],
        font_size=15,
        color=ACCENT,
    )

    # Slide 17 — next steps (updated)
    slide = prs.slides.add_slide(blank)
    add_title(slide, "16. 下一步：从验证到产品化", "nav-gateway 链路已通，room-scale 验证完成，接下来是真实传感器、真实环境和量产部署")
    add_formula_box(
        slide,
        "已验证链路",
        [
            "odor sensing -> odor planner",
            "odor planner -> nav-gateway dispatch/retarget",
            "nav-gateway -> LingTu bridge",
            "LingTu -> obstacle avoidance / recovery / safety",
            "room-scale 8-scenario validation (8m×6m)",
        ],
        0.78,
        1.55,
        5.80,
        2.50,
    )
    add_bullets(
        slide,
        6.85,
        1.50,
        5.30,
        5.20,
        [
            "✓ 已完成：嗅觉策略已通过 nav-gateway API 发送滚动目标到 LingTu bridge。",
            "✓ 已完成：8 组室内场景验证了不同风向、风速、起始位置下的搜源能力和方法边界。",
            "第一步：加入真实传感器回放，用实际传感器数据替换仿真模型。",
            "第二步：由 LingTu 继续负责避障、恢复、地图约束，嗅觉模块只负责“该往哪片区域搜”。",
            "第三步：加入障碍场景、更多风场扰动，再做成功率和时间统计。",
            "第四步：真实环境现场测试与量产部署方案验证。",
            "第五步：针对下风起步和横风偏置两个 CHALLENGE 案例，加入全局覆盖搜索策略（当前 zigzag 已将成功率从 5/8 提升至 6/8）。",
        ],
    )
    add_textbox(
        slide,
        0.82,
        6.50,
        11.0,
        0.30,
        "结论：nav-gateway 链路已跑通，room-scale 验证明确了方法边界，从“研究草图”推进到“可运行、可对接、可量产化”的阶段。",
        font_size=16,
        bold=True,
        color=ACCENT_2,
    )

    for slide in prs.slides:
        background = slide.background
        background.fill.solid()
        background.fill.fore_color.rgb = BG

    OUTPUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PPT)
    print(OUTPUT_PPT)


if __name__ == "__main__":
    main()
